import streamlit as st
from pptx import Presentation
from llm_client import query_together
import json
import re
from report_generator import generate_pdf_report
import pandas as pd

REQUIRED_TYPES = [
    "Problem", "Solution", "Market", "Business Model",
    "Competition", "Team", "Financials", "Traction", "Funding Ask"
]

def count_slides(file):
    prs = Presentation(file)
    return len(prs.slides)

def extract_slide_text(file):
    prs = Presentation(file)
    slide_texts = []
    for i, slide in enumerate(prs.slides):
        content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    content.append(text)
        slide_texts.append({"slide_num": i + 1, "text": " ".join(content)})
    return slide_texts

def build_classification_prompt(slides):
    intro = (
        "You are an AI assistant that classifies slides from a startup pitch deck. "
        "Each slide must be assigned one of the following categories:\n"
        "Problem, Solution, Market, Business Model, Competition, Team, Financials, Traction, Funding Ask, Unclassified.\n"
        "For each slide, respond in this JSON format:\n"
        "{\"slide_num\": 1, \"category\": \"Team\"}\n"
        "Here are the slides:\n"
    )
    body = ""
    for slide in slides:
        body += f"\nSlide {slide['slide_num']}:\n{slide['text']}\n"
    return intro + body

def classify_slides_with_llm(slides):
    prompt = build_classification_prompt(slides)
    raw_output = query_together(prompt)
    return raw_output

def parse_classification_output(output_text):
    try:
        return json.loads(output_text)
    except json.JSONDecodeError as e:
        raise ValueError("LLM returned invalid JSON") from e

def classify_single_slide(text):
    prompt = (
        "Task: Classify the following startup pitch deck slide into exactly one category.\n\n"
        "Categories:\n"
        "1. Problem - Market problems or pain points\n"
        "2. Solution - Product or service solutions\n"
        "3. Market - Market size, opportunity, target market\n"
        "4. Business Model - Revenue model, pricing, strategy\n"
        "5. Competition - Competitors or competitive advantage\n"
        "6. Team - Team members, experience, expertise\n"
        "7. Financials - Financial projections, metrics, funding\n"
        "8. Traction - Growth, customers, achievements\n"
        "9. Funding Ask - Funding needs, investment terms\n\n"
        f"Slide Content:\n{text}\n\n"
        "Instructions:\n"
        "1. Choose the single most appropriate category\n"
        "2. Return ONLY the category name\n"
        "3. Do not include any explanations or additional text\n"
        "4. Use exact category names as listed above"
    )
    try:
        response = query_together(prompt, max_tokens=50)
        # Clean and validate the response
        category = response.strip().split()[0].capitalize()  # Take first word and capitalize
        if category in REQUIRED_TYPES:
            return category
        print(f"Invalid category returned: {category}")  # Debug logging
        return "Unclassified"
    except Exception as e:
        print(f"Error classifying slide: {e}")
        return "Unclassified"

def classify_all_slides(slides):
    results = []
    for slide in slides:
        try:
            label = classify_single_slide(slide["text"])
            slide["category"] = label
            print(f"Classified slide {slide['slide_num']} as: {label}")  # Debug logging
        except Exception as e:
            print(f"Error processing slide {slide['slide_num']}: {e}")
            slide["category"] = "Unclassified"
        results.append(slide)
    return results
    
    
def build_analysis_prompt(slides):
    classified = [s for s in slides if s.get("category") != "Unclassified"]

    header = (
        "You are an AI investment analyst. Based on the following pitch deck content, produce an investment thesis. "
        "IMPORTANT: Return your response as a VALID JSON object with the following fields:\n\n"
        "- recommendation: one of [\"Strong Buy\", \"Hold\", \"Pass\"]\n"
        "- overall_score: integer 0–100\n"
        "- processing_date: use current UTC in format DD-MM-YYYY HH:MM:SS UTC\n"
        "- confidence_score: integer 0–100\n"
        "- strengths: list of 3–5 strings\n"
        "- weaknesses: list of 3–5 strings\n"
        "- recommendations: string (100–200 words)\n"
        "- categories: list of 9 objects, each with:\n"
        "  - name (category name),\n"
        "  - score (0–10),\n"
        "  - weight (int %),\n"
        "  - feedback (50–150 words)\n\n"
        "Use these fixed weights: Problem 10, Solution 15, Market 20, Business Model 15, Competition 10, Team 15, "
        "Traction 10, Financials 10, Clarity 5\n\n"
        "STRICT JSON FORMATTING RULES:\n"
        "1. Return ONLY a valid JSON object, nothing else\n"
        "2. Use double quotes for all keys and string values\n"
        "3. Do not include any explanatory text before or after the JSON\n"
        "4. Do not use markdown code blocks or formatting\n"
        "5. Ensure all strings are properly escaped\n\n"
        "Classified Slides:\n"
    )

    grouped = {}
    for slide in classified:
        cat = slide["category"]
        grouped.setdefault(cat, []).append(slide["text"])

    body = ""
    for cat in grouped:
        content = "\n".join(grouped[cat])
        body += f"\n---\nCategory: {cat}\n{content}\n"

    return header + body

def analyze_pitch(slides):
    prompt = build_analysis_prompt(slides)
    
    # Add a specific instruction to ensure valid JSON output
    json_instruction = (
        "\n\nIMPORTANT: Your response MUST be a valid JSON object. "
        "Do not include any text, markdown formatting, or code blocks outside the JSON object. "
        "Ensure all string values use double quotes and are properly escaped."
    )
    
    enhanced_prompt = prompt + json_instruction
    
    # Request more tokens to ensure complete response
    raw_output = query_together(enhanced_prompt, max_tokens=3500)
    
    # Basic validation check
    if not raw_output.strip().startswith('{') and not '{' in raw_output:
        print("Warning: LLM response does not appear to contain JSON")
    
    return raw_output

def parse_analysis_output(raw_text):
    # Clean up the raw text to extract valid JSON
    cleaned = raw_text.strip()
    
    # Remove any markdown code block indicators
    cleaned = re.sub(r'^```json\s*|\s*```$', '', cleaned, flags=re.MULTILINE)
    cleaned = re.sub(r'^```\s*|\s*```$', '', cleaned, flags=re.MULTILINE)
    
    # Remove any 'json' text at the beginning
    cleaned = re.sub(r'^json\s*', '', cleaned, flags=re.MULTILINE)
    
    # Remove any leading/trailing whitespace on each line
    cleaned = re.sub(r'^\s*|\s*$', '', cleaned, flags=re.MULTILINE)
    
    # Remove any markdown section separators
    cleaned = re.sub(r'^---\s*', '', cleaned, flags=re.MULTILINE)
    
    # Try to find JSON object boundaries if there's text before or after
    match = re.search(r'\{.*\}', cleaned, re.DOTALL)
    if match:
        cleaned = match.group(0)
    
    print("Cleaned JSON:", cleaned[:100], "..." if len(cleaned) > 100 else "")
    
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError as e:
        st.error(f"Invalid JSON: {str(e)}")
        raise ValueError("Invalid JSON in LLM analysis output.") from e



st.set_page_config(page_title="Investment Thesis Generator", layout="centered")

st.title("Automated Investment Thesis Generator (POC)")

st.markdown("""
Upload a startup pitch deck in .pptx format. The deck should have 5–20 slides and include at least 3 of the following:  
- Problem  
- Solution/Product  
- Market  
- Business Model  
- Competition  
- Team  
- Financials  
- Traction  
- Funding Ask
""")

uploaded_file = st.file_uploader("Upload your PowerPoint (.pptx only)", type=["pptx"])

if uploaded_file:
    if uploaded_file.size > 50 * 1024 * 1024:
        st.error("File size exceeds 50MB.")
    else:
        slide_count = count_slides(uploaded_file)
        if not (5 <= slide_count <= 20):
            st.error(f"Invalid slide count: {slide_count}. Must be between 5 and 20.")
        else:
            st.success(f"Valid file with {slide_count} slides.")

            if "slide_data" not in st.session_state:
                if st.button("Proceed to Analysis"):
                    with st.spinner("Extracting text from slides..."):
                        slide_data = extract_slide_text(uploaded_file)
                        st.session_state["slide_data"] = slide_data
                        st.success("Text extraction complete.")

if "slide_data" in st.session_state:
    slide_data = st.session_state["slide_data"]

    if st.button("Classify Slides with AI"):
        with st.spinner("Classifying each slide..."):
            classified_slides = classify_all_slides(slide_data)
            print([s["category"] for s in classified_slides])
            st.session_state["classified_slides"] = classified_slides

            found_types = {s["category"] for s in classified_slides if s["category"] in REQUIRED_TYPES}
            if len(found_types) < 3:
                st.error(f"Only {len(found_types)} valid categories found: {', '.join(found_types)}")
            else:
                st.success(f"Classification complete. Categories found: {', '.join(found_types)}")

if "classified_slides" in st.session_state:
    if st.button("Run Investment Analysis"):
        with st.spinner("Analyzing pitch deck..."):
            raw_analysis = None
            try:
                raw_analysis = analyze_pitch(st.session_state["classified_slides"])
                # Always show raw output for debugging
                with st.expander("Show raw LLM output (JSON string)"):
                    st.code(raw_analysis, language="json")
                parsed = parse_analysis_output(raw_analysis)
                st.session_state["analysis_result"] = parsed
                st.success("Analysis complete.")
            except Exception as e:
                st.error(str(e))
                if raw_analysis:
                    st.markdown("#### Raw LLM output (for debugging)")
                    st.code(raw_analysis, language="json")
                st.stop()

        # Optional debug: display result
        st.subheader("Investment Thesis Summary")
        col1, col2, col3 = st.columns(3)
        col1.metric("Recommendation", parsed.get("recommendation", "-"))
        col2.metric("Overall Score", parsed.get("overall_score", "-"))
        col3.metric("Confidence Score", parsed.get("confidence_score", "-"))

        st.markdown("---")
        st.markdown("### Strengths")
        for s in parsed.get("strengths", []):
            st.markdown(f"- {s}")

        st.markdown("### Weaknesses")
        for w in parsed.get("weaknesses", []):
            st.markdown(f"- {w}")

        st.markdown("### Recommendations")
        st.markdown(parsed.get("recommendations", "-"))

        st.markdown("---")
        st.markdown("### Category-wise Analysis")
        cat_data = parsed.get("categories", [])
        if cat_data:
            # Prepare data for DataFrame
            df = pd.DataFrame([
                {
                    "Category": c.get("name", "-"),
                    "Score": c.get("score", "-"),
                    "Weight (%)": c.get("weight", "-"),
                    "Feedback": c.get("feedback", "-")
                }
                for c in cat_data
            ])
            # Show table with expanders for feedback
            for idx, row in df.iterrows():
                with st.expander(f"{row['Category']} (Score: {row['Score']}, Weight: {row['Weight (%)']}%)"):
                    st.markdown(row['Feedback'])
        else:
            st.warning("No category analysis available.")

        st.markdown("---")
        if st.button("Generate PDF Report"):
            pdf_bytes, pdf_name = generate_pdf_report(parsed, startup_name="DemoStartup")
            st.download_button(label="Download Report", data=pdf_bytes, file_name=pdf_name, mime="application/pdf")